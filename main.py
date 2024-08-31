import tkinter as tk
from tkinter import filedialog, messagebox
from PIL import Image, ImageTk
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Function to create the presentation
def create_presentation():
    prs = Presentation()

    def add_background(slide, image_path):
        # Load and resize the image
        image = Image.open(image_path)
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        image = image.resize((slide_width, slide_height), Image.ANTIALIAS)
        image_path_resized = "resized_background.jpg"
        image.save(image_path_resized)
        
        # Add the resized image to the slide
        slide.shapes.add_picture(image_path_resized, 0, 0, width=slide_width, height=slide_height)

    def add_title_slide(prs, title_text, subtitle_text, image_path):
        slide = prs.slides.add_slide(prs.slide_layouts[0])

        # Add background image
        add_background(slide, image_path)

        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = title_text
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text

        subtitle.text = subtitle_text
        subtitle.text_frame.paragraphs[0].font.size = Pt(24)
        subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text

    def add_agenda_slide(prs, title_text, agenda_items, image_path):
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        # Add background image
        add_background(slide, image_path)

        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text

        content = slide.placeholders[1].text_frame
        for item in agenda_items:
            p = content.add_paragraph()
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(255, 255, 255)  # White text

    def add_content_slide(prs, title_text, content_text, image_path):
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        # Add background image
        add_background(slide, image_path)

        title = slide.shapes.title
        title.text = title_text
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text

        # Add content
        textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(5))
        text_frame = textbox.text_frame
        p = text_frame.add_paragraph()
        p.text = content_text
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255, 255, 255)  # White text
        text_frame.word_wrap = True

    # Gather data from the UI
    title_text = title_entry.get()
    subtitle_text = subtitle_entry.get()
    agenda_items = agenda_text.get("1.0", tk.END).strip().split("\n")
    image_path = filedialog.askopenfilename(title="Select Background Image", filetypes=[("Image Files", "*.jpg;*.jpeg;*.png")])

    if not image_path:
        messagebox.showerror("Error", "No background image selected.")
        return

    # Create the presentation
    add_title_slide(prs, title_text, subtitle_text, image_path)
    add_agenda_slide(prs, "Agenda", agenda_items, image_path)

    content_sections = [
        ("Introduction", "Introduce your topic and give an overview of what will be covered."),
        ("Code Overview", "Highlight key parts of your code. Explain the logic behind them."),
        ("Functionality Demonstration", "Show a live demo or walk through how your application works."),
        ("Optimization Techniques", "Discuss how you optimized your code, including error handling and modularity."),
        ("Audience Engagement", "Engage the audience with polls, questions, and interactive elements."),
        ("Conclusion", "Summarize the main points and provide a call to action.")
    ]

    for section in content_sections:
        add_content_slide(prs, section[0], section[1], image_path)

    # Save the presentation with error handling
    try:
        save_path = filedialog.asksaveasfilename(defaultextension=".pptx", initialfile="suraj.pptx", filetypes=[("PowerPoint Files", "*.pptx")])
        if save_path:
            prs.save(save_path)
            messagebox.showinfo("Success", f"Presentation saved as {save_path}")
        else:
            messagebox.showwarning("Warning", "File not saved. No file path provided.")
    except Exception as e:
        messagebox.showerror("Error", f"Failed to save the presentation: {str(e)}")

# Tkinter UI setup
root = tk.Tk()
root.title("Presentation Creator")
root.geometry("600x400")

# Load background image
bg_image = Image.open("ppt.jpg")
bg_image = bg_image.resize((600, 400), Image.LANCZOS)  # Resize the image to fit the window
bg_photo = ImageTk.PhotoImage(bg_image)

# Create a label to hold the background image
bg_label = tk.Label(root, image=bg_photo)
bg_label.place(x=0, y=0, relwidth=1, relheight=1)  # Set the label to cover the entire window

# Title input
tk.Label(root, text="Presentation Title:", font=("Arial", 14), bg='#007acc', fg='white').pack(pady=10)
title_entry = tk.Entry(root, width=40, font=("Arial", 12))
title_entry.pack(pady=5)

# Subtitle input
tk.Label(root, text="Presentation Subtitle:", font=("Arial", 14), bg='#007acc', fg='white').pack(pady=10)
subtitle_entry = tk.Entry(root, width=40, font=("Arial", 12))
subtitle_entry.pack(pady=5)

# Agenda input
tk.Label(root, text="Agenda (one item per line):", font=("Arial", 14), bg='#007acc', fg='white').pack(pady=10)
agenda_text = tk.Text(root, height=6, width=40, font=("Arial", 12))
agenda_text.pack(pady=5)

# Create Presentation button
create_btn = tk.Button(root, text="Create Presentation", font=("Arial", 14), command=create_presentation, bg='#007acc', fg='white')
create_btn.pack(pady=20)

root.mainloop()
